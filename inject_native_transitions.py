# -*- coding: utf-8 -*-
"""
inject_native_transitions.py
使用 python-pptx / lxml 为 PPTX 中的每一页注入标准 OpenXML 原生淡入淡出切换动画 (Smooth Fade Transition)
这能 100% 保证在 WPS Office 和 Microsoft PowerPoint 播放时拥有极度丝滑平柔的换页体验！
"""

import os
import sys

sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.oxml import parse_xml

pptx_path = r"D:\jinyuliangyuan\宣传用PPT\今遇莨缘_香云纱品牌宣传_经典画风版.pptx"

if not os.path.exists(pptx_path):
    print("❌ 找不到文件：", pptx_path)
    sys.exit(1)

prs = Presentation(pptx_path)
print(f"📄 正在为 {len(prs.slides)} 页幻灯片注入原生 OpenXML 丝滑淡入淡出切换动画...")

# XML 标准淡入淡出 (Smooth Fade, Slow speed)
# spd="slow" 代表慢速柔和，<p:fade thresh="80"/> 或 <p:fade/>
fade_xml_str = """
<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow">
    <p:fade/>
</p:transition>
"""

count = 0
for i, slide in enumerate(prs.slides):
    # 获取幻灯片根节点 <p:sld>
    sld_elem = slide.element
    
    # 检查是否已存在 transition
    existing_trans = sld_elem.find('{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    if existing_trans is not None:
        sld_elem.remove(existing_trans)
        
    # 解析新 XML 并追加至 <p:sld>
    new_trans_elem = parse_xml(fade_xml_str)
    sld_elem.append(new_trans_elem)
    count += 1
    print(f"   ✅ Slide {i+1}: 成功注入原生 Fade 柔和平滑切换")

prs.save(pptx_path)
print(f"\n🎉 保存成功！{count} 页已全部应用顶级丝滑切换动画！")
