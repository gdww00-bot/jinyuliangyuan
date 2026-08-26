# -*- coding: utf-8 -*-
import sys
import os

# Force UTF-8 output
sys.stdout.reconfigure(encoding='utf-8')

try:
    from pptx import Presentation
except ImportError:
    os.system(r"C:\Users\Administrator\AppData\Local\Programs\Python\Python313\python.exe -m pip install python-pptx -q")
    from pptx import Presentation

pptx_dir = r"D:\jinyuliangyuan\宣传用PPT"
files = [
    "今遇莨缘_香云纱品牌宣传.pptx",
    "今遇莨缘香云纱品牌宣传PPT.pptx",
    "今遇莨缘香云纱品牌宣传PPT (1).pptx",
]

for fname in files:
    fpath = os.path.join(pptx_dir, fname)
    print("=" * 60)
    print(f"[FILE] {fname}")
    print(f"   Size: {os.path.getsize(fpath)/1024:.0f} KB")
    try:
        prs = Presentation(fpath)
        print(f"   Slides: {len(prs.slides)}")
        print(f"   Dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
        for i, slide in enumerate(prs.slides):
            texts = []
            img_count = 0
            shape_types = []
            for shape in slide.shapes:
                shape_types.append(shape.shape_type)
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    img_count += 1
            layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
            print(f"\n  [Slide {i+1}] Layout: {layout_name} | Images: {img_count}")
            for t in texts[:10]:
                print(f"    > {t[:70]}")
            if len(texts) > 10:
                print(f"    ... +{len(texts)-10} more text items")
    except Exception as e:
        print(f"   ERROR: {e}")
    print()
